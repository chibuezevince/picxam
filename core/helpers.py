import os

import docx
import fitz
from pptx import Presentation

from core.extract_images import extract_office, extract_pdf
from core.models import Setting
from django.conf import settings
from pathlib import Path
from typing import List


def get_setting(key, default=None):
    try:
        return Setting.objects.get(key=key).value
    except Setting.DoesNotExist:
        return default


def upload_file(f, folder=None):
    subdir = os.path.join("uploads", folder) if folder else "uploads"
    dest = os.path.join(settings.MEDIA_ROOT, subdir)
    os.makedirs(dest, exist_ok=True)
    file_path = os.path.join(dest, f.name)

    with open(file_path, "wb+") as destination:
        for chunk in f.chunks():
            destination.write(chunk)
    rel_path = os.path.join(subdir, f.name)

    return rel_path


def extract_images(
    file_path: str,
    output_dir: str | None = None,
) -> List[str]:
    """
    Extracts images from a single PDF, PPTX, or DOCX file.
    Filters out images smaller than min_width/min_height.

    Returns:
        List of absolute paths to the saved image files.
    """
    file_path = Path(file_path).resolve()
    image_directory = (
        settings.MEDIA_ROOT / "extracted_images"
        if not output_dir
        else settings.MEDIA_ROOT / "extracted_images" / str(output_dir)
    )

    output_dir = Path(image_directory)
    output_dir.mkdir(parents=True, exist_ok=True)

    extension = Path(file_path).suffix.lstrip(".")

    extracted_images = []
    match extension:
        case "pptx" | "docx":
            extracted_images = extract_office(
                file_path,
                output_dir,
                extension,
            )

        case "pdf":
            extracted_images = extract_pdf(file_path, output_dir)

    return extracted_images


def get_relative_path(absolute_path: str) -> str:
    """Convert an absolute path to one relative to MEDIA_ROOT."""
    return os.path.relpath(absolute_path, settings.MEDIA_ROOT)


def extract_text(abs_path: str) -> str:
    """Extracts text from PDF, PPTX, or DOCX files."""
    path = Path(abs_path)
    extension = path.suffix.lstrip(".").lower()

    match extension:
        case "pdf":
            text_parts = []
            with fitz.open(abs_path) as doc:
                for page in doc:
                    text_parts.append(page.get_text())
            return "\n".join(text_parts)

        case "pptx":
            text_parts = []
            prs = Presentation(abs_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text_parts.append(paragraph.text)
            return "\n".join(text_parts)

        case "docx":
            doc = docx.Document(abs_path)
            return "\n".join(paragraph.text for paragraph in doc.paragraphs)

        case _:
            return ""
